"""
报告生成工具 - 从分析结果到交付物

支持生成 PDF、PPT、Markdown 格式的分析报告。
实现"一键出报告"亮点。
"""

import os
import json
from datetime import datetime
from typing import Optional, Dict, Any, List
from hello_agents.tools.base import Tool, ToolParameter
from hello_agents.tools.response import ToolResponse
from hello_agents.tools.errors import ToolErrorCode


class ReportGeneratorTool(Tool):
    """报告生成工具 - 从数据到交付物"""

    def __init__(self, output_dir: str = "./output/reports"):
        """
        初始化报告生成工具

        Args:
            output_dir: 报告输出目录
        """
        super().__init__(
            name="report_generator",
            description="生成分析报告，支持PDF、PPT、Markdown格式。"
                        "输入分析结论和图表路径，输出完整的分析报告文件。"
        )
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="report_type",
                type="string",
                description="报告类型: pdf/ppt/markdown",
                required=True
            ),
            ToolParameter(
                name="title",
                type="string",
                description="报告标题",
                required=True
            ),
            ToolParameter(
                name="content",
                type="string",
                description="报告内容（JSON格式，包含sections数组）",
                required=True
            ),
            ToolParameter(
                name="chart_paths",
                type="string",
                description="图表文件路径列表（JSON数组格式）",
                required=False
            ),
            ToolParameter(
                name="filename",
                type="string",
                description="输出文件名（不含扩展名）",
                required=False
            )
        ]

    def run(self, parameters: dict) -> ToolResponse:
        """
        生成报告

        Args:
            parameters: 报告参数

        Returns:
            ToolResponse 包含报告文件路径
        """
        report_type = parameters.get("report_type", "markdown")
        title = parameters.get("title", "数据分析报告")
        content_str = parameters.get("content", "{}")
        chart_paths_str = parameters.get("chart_paths", "[]")
        filename = parameters.get("filename", "")

        try:
            # 解析内容
            if isinstance(content_str, str):
                content = json.loads(content_str)
            else:
                content = content_str

            # 解析图表路径
            if isinstance(chart_paths_str, str):
                chart_paths = json.loads(chart_paths_str)
            else:
                chart_paths = chart_paths_str

            # 生成文件名
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"report_{timestamp}"

            # 根据类型生成
            if report_type == "markdown":
                filepath = self._generate_markdown(title, content, chart_paths, filename)
            elif report_type == "pdf":
                filepath = self._generate_pdf(title, content, chart_paths, filename)
            elif report_type == "ppt":
                filepath = self._generate_ppt(title, content, chart_paths, filename)
            else:
                return ToolResponse.error(
                    code=ToolErrorCode.INVALID_PARAM,
                    message=f"不支持的报告类型: {report_type}"
                )

            return ToolResponse.success(
                text=f"报告已生成: {filepath}",
                data={
                    "filepath": filepath,
                    "report_type": report_type,
                    "title": title
                }
            )

        except json.JSONDecodeError as e:
            return ToolResponse.error(
                code=ToolErrorCode.INVALID_PARAM,
                message=f"内容JSON解析失败: {str(e)}"
            )
        except Exception as e:
            return ToolResponse.error(
                code=ToolErrorCode.EXECUTION_ERROR,
                message=f"报告生成失败: {str(e)}"
            )

    def _generate_markdown(self, title: str, content: Dict, chart_paths: List, filename: str) -> str:
        """生成 Markdown 报告"""
        filepath = os.path.join(self.output_dir, f"{filename}.md")

        lines = [
            f"# {title}",
            f"",
            f"> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"",
            "---",
            ""
        ]

        # 添加各个章节
        sections = content.get("sections", [])
        for section in sections:
            section_title = section.get("title", "")
            section_content = section.get("content", "")

            if section_title:
                lines.append(f"## {section_title}")
                lines.append("")

            if section_content:
                lines.append(section_content)
                lines.append("")

        # 添加图表
        if chart_paths:
            lines.append("## 数据可视化")
            lines.append("")
            for i, chart_path in enumerate(chart_paths):
                if os.path.exists(chart_path):
                    lines.append(f"![图表{i+1}]({chart_path})")
                    lines.append("")

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))

        return filepath

    def _generate_pdf(self, title: str, content: Dict, chart_paths: List, filename: str) -> str:
        """生成 PDF 报告"""
        try:
            from fpdf import FPDF
        except ImportError:
            # 回退到 Markdown
            return self._generate_markdown(title, content, chart_paths, filename)

        filepath = os.path.join(self.output_dir, f"{filename}.pdf")

        pdf = FPDF()
        pdf.add_page()

        # 尝试添加中文字体
        try:
            pdf.add_font('SimHei', '', 'C:/Windows/Fonts/simhei.ttf', uni=True)
            pdf.set_font('SimHei', size=16)
        except Exception:
            pdf.set_font('Arial', size=16)

        # 标题
        pdf.cell(0, 15, title, ln=True, align='C')
        pdf.ln(5)

        # 生成时间
        try:
            pdf.set_font('SimHei', size=10)
        except Exception:
            pdf.set_font('Arial', size=10)
        pdf.cell(0, 10, f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", ln=True, align='C')
        pdf.ln(10)

        # 章节内容
        sections = content.get("sections", [])
        for section in sections:
            section_title = section.get("title", "")
            section_content = section.get("content", "")

            if section_title:
                try:
                    pdf.set_font('SimHei', size=14)
                except Exception:
                    pdf.set_font('Arial', size=14)
                pdf.cell(0, 10, section_title, ln=True)
                pdf.ln(3)

            if section_content:
                try:
                    pdf.set_font('SimHei', size=11)
                except Exception:
                    pdf.set_font('Arial', size=11)
                pdf.multi_cell(0, 8, section_content)
                pdf.ln(5)

        # 插入图表
        for chart_path in chart_paths:
            if os.path.exists(chart_path):
                try:
                    pdf.image(chart_path, w=170)
                    pdf.ln(5)
                except Exception:
                    pass

        pdf.output(filepath)
        return filepath

    def _generate_ppt(self, title: str, content: Dict, chart_paths: List, filename: str) -> str:
        """生成 PPT 报告"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            # 回退到 Markdown
            return self._generate_markdown(title, content, chart_paths, filename)

        filepath = os.path.join(self.output_dir, f"{filename}.pptx")

        prs = Presentation()

        # 标题页
        slide_layout = prs.slide_layouts[0]  # 标题页布局
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # 内容页
        sections = content.get("sections", [])
        for section in sections:
            section_title = section.get("title", "分析结果")
            section_content = section.get("content", "")

            slide_layout = prs.slide_layouts[1]  # 标题+内容布局
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section_title

            if section_content:
                slide.placeholders[1].text = section_content

        # 图表页
        for chart_path in chart_paths:
            if os.path.exists(chart_path):
                slide_layout = prs.slide_layouts[5]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(chart_path, Inches(1), Inches(1), Inches(8), Inches(5))

        prs.save(filepath)
        return filepath
